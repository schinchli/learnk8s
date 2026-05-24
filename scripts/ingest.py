#!/usr/bin/env python3
"""
EKS Course Knowledge Ingestion Pipeline — full learning path
Sources:
  1. slides       PPTX instructor decks (modules 01–09)
  2. notes        course-flashcards/module-XX/notes.md (modules 02–09)
  3. flashcards   course-flashcards/module-XX/index.html (modules 01–09)
  4. reference    lab-1-deploying-pods/ HTML reference pages
  5. lab_steps    animated-lab-*.html narrations (engine-based labs)
  6. lab_notes    lab-1-deploying-pods/notes/*.md

Usage:
  python3 scripts/ingest.py                      # ingest all sources
  python3 scripts/ingest.py --source slides
  python3 scripts/ingest.py --source reference
  python3 scripts/ingest.py --module 01
  python3 scripts/ingest.py --clear              # wipe table first
  python3 scripts/ingest.py --clear --dry-run    # count only

Requirements: pip3 install openai python-pptx supabase
"""

import os, sys, json, re, time, argparse, hashlib
from pathlib import Path
from html.parser import HTMLParser
from pptx import Presentation
from openai import OpenAI
from supabase import create_client

CORPUS = 'eks-coreks'

# ── Config ────────────────────────────────────────────────────────────
REPO_ROOT   = Path(__file__).parent.parent
PPTX_BASE   = Path('/Users/schinchli/Downloads/200-COREKS-22-EN-PPTX.2.2.3-20260417145107 2/')
FLASHCARDS  = REPO_ROOT / 'course-flashcards'
LAB_DIR     = REPO_ROOT / 'lab-1-deploying-pods'
EMBED_MODEL = 'text-embedding-3-small'
EMBED_DIMS  = 1536

MODULE_TITLES = {
    '01': 'Kubernetes Fundamentals',
    '02': 'Amazon EKS Fundamentals',
    '03': 'Building and Maintaining an Amazon EKS Cluster',
    '04': 'Deploying Applications to Your Amazon EKS Cluster',
    '05': 'Managing Applications at Scale in Amazon EKS',
    '06': 'Managing Networking in Amazon EKS',
    '07': 'Configuring Observability in an Amazon EKS Cluster',
    '08': 'Managing Storage in Amazon EKS',
    '09': 'Managing Security in Amazon EKS',
}

PPTX_FILES = {
    '01': '01_KubernetesFundamentals_InstructorDeck.pptx',
    '02': '02_AmazonEKSFundamentals_InstructorDeck.pptx',
    '03': '03_BuildingandMaintaininganAmazonEKSCluster_InstructorDeck.pptx',
    '04': '04_DeployingApplicationstoyourEKSCluster_InstructorDeck.pptx',
    '05': '05_ManagingApplicationsatScaleinAmazonEKS_InstructorDeck.pptx',
    '06': '06_ManagingNetworkinginAmazonEKS_InstructorDeck.pptx',
    '07': '07_ConfiguringObservabilityinanAmazonEKSCluster_InstructorDeck.pptx',
    '08': '08_ManagingStorageinAmazonEKS_InstructorDeck.pptx',
    '09': '09_ManagingSecurityinAmazonEKS_InstructorDeck.pptx',
}

# Animated labs that use the shared engine (have narration:'...' fields)
ENGINE_LABS = {
    '03': 'animated-lab-building-eks-cluster.html',
    '04': 'animated-lab-deploying-apps.html',
    '05': 'animated-lab-scale-gitops.html',
    '06': 'animated-lab-networking.html',
    '07': 'animated-lab-observability.html',
    '08': 'animated-lab-storage.html',
}

# Reference learning pages → (file, module, page_title)
REFERENCE_PAGES = [
    ('learn-visually.html',    '01', 'Visual Reference — Kubernetes'),
    ('deep-dive.html',         '01', 'Deep Dive — Every K8s Component'),
    ('cluster-explained.html', '01', 'Live Cluster Visual — Every Component Explained'),
    ('visual-concepts.html',   '01', 'Visual Concepts — Animated K8s'),
]

# ── Clients ───────────────────────────────────────────────────────────
def get_clients():
    url  = os.environ.get('SUPABASE_URL')
    key  = os.environ.get('SUPABASE_SERVICE_ROLE_KEY')
    okey = os.environ.get('OPENAI_API_KEY')
    if not url:  raise RuntimeError('SUPABASE_URL not set')
    if not key:  raise RuntimeError('SUPABASE_SERVICE_ROLE_KEY not set')
    if not okey: raise RuntimeError('OPENAI_API_KEY not set')
    return OpenAI(api_key=okey), create_client(url, key)

# ── Embedding ─────────────────────────────────────────────────────────
def embed_batch(oai, texts: list[str]) -> list[list[float]]:
    results = []
    for i in range(0, len(texts), 50):
        batch = [t[:8000] for t in texts[i:i+50]]
        for attempt in range(3):
            try:
                resp = oai.embeddings.create(model=EMBED_MODEL, input=batch)
                results.extend([d.embedding for d in resp.data])
                time.sleep(0.25)
                break
            except Exception as e:
                if attempt == 2: raise
                print(f'    retry embed ({e})…')
                time.sleep(2 ** attempt)
    return results

def upsert_chunks(sb, chunks: list[dict], oai, dry_run=False) -> int:
    """Writes into the generic knowledge_chunks table with corpus='eks-coreks'.
    Idempotent — the (corpus, content_hash) unique index causes duplicate
    rows to be ignored. Module/module_title/topic move into metadata jsonb."""
    if not chunks:
        return 0
    if dry_run:
        print(f'    [dry-run] would embed+insert {len(chunks)} chunks')
        return len(chunks)
    texts = [c['content'] for c in chunks]
    print(f'    embedding {len(texts)} chunks…')
    embeddings = embed_batch(oai, texts)
    rows = []
    for chunk, emb in zip(chunks, embeddings):
        content = chunk['content'][:4000]
        chash = hashlib.sha256(content.encode('utf-8')).hexdigest()
        metadata = dict(chunk.get('metadata') or {})
        metadata.update({
            'module':       chunk['module'],
            'module_title': chunk.get('module_title', ''),
            'topic':        chunk.get('topic'),
        })
        rows.append({
            'corpus':       CORPUS,
            'source_type':  chunk['source_type'],
            'content_hash': chash,
            'title':        (chunk.get('title') or '')[:500],
            'content':      content,
            'embedding':    emb,
            'metadata':     metadata,
        })
    inserted = 0
    for i in range(0, len(rows), 50):
        # upsert on (corpus, content_hash) — skips duplicates silently
        sb.table('knowledge_chunks') \
          .upsert(rows[i:i+50], on_conflict='corpus,content_hash', ignore_duplicates=True) \
          .execute()
        inserted += len(rows[i:i+50])
    return inserted

# ══════════════════════════════════════════════════════════════════════
# Source 1: PPTX slides
# ══════════════════════════════════════════════════════════════════════
def ingest_slides(oai, sb, only_module=None, dry_run=False) -> int:
    total = 0
    modules = {only_module: PPTX_FILES[only_module]} if only_module and only_module in PPTX_FILES else PPTX_FILES
    for mod, fname in modules.items():
        path = PPTX_BASE / fname
        if not path.exists():
            print(f'  [skip] {fname} — not found'); continue
        prs = Presentation(str(path))
        mod_title = MODULE_TITLES.get(mod, '')
        chunks, current_topic = [], mod_title
        for slide_idx, slide in enumerate(prs.slides):
            texts = [s.text.strip() for s in slide.shapes if hasattr(s,'text') and s.text.strip()]
            if not texts: continue
            if len(texts) <= 2 and all(len(t) < 80 for t in texts):
                candidate = texts[0].replace('\n',' ').strip()
                if candidate and 'running containers' not in candidate.lower():
                    current_topic = candidate
                continue
            slide_num = slide_idx + 1
            title = texts[0].replace('\x0b',' ').replace('\n',' ').strip()
            body  = '\n'.join(texts[1:]).replace('\x0b','\n')
            if 'knowledge check' in title.lower(): current_topic = 'Knowledge Check'
            elif 'module summary' in title.lower(): current_topic = 'Module Summary'
            elif 'lab ' in title.lower():           current_topic = 'Lab Exercise'
            chunks.append({
                'module': mod, 'module_title': mod_title,
                'source_type': 'slide', 'topic': current_topic,
                'title': title[:500], 'content': f"{title}\n\n{body}".strip(),
                'metadata': {'slide': slide_num, 'pptx': fname},
            })
        n = upsert_chunks(sb, chunks, oai, dry_run)
        print(f'  M{mod} slides: {n} chunks')
        total += n
    return total

# ══════════════════════════════════════════════════════════════════════
# Source 2: notes.md study guides
# ══════════════════════════════════════════════════════════════════════
def ingest_notes(oai, sb, only_module=None, dry_run=False) -> int:
    total = 0
    all_modules = ['01','02','03','04','05','06','07','08','09']
    modules = [only_module] if only_module else all_modules
    for mod in modules:
        path = FLASHCARDS / f'module-{mod}' / 'notes.md'
        if not path.exists(): continue
        text = path.read_text(encoding='utf-8')
        mod_title = MODULE_TITLES.get(mod, '')
        chunks = []
        for section in re.split(r'\n(?=## )', text):
            section = section.strip()
            if not section or len(section) < 60: continue
            lines = section.split('\n')
            raw_title = lines[0].lstrip('#').strip()
            body = '\n'.join(lines[1:]).strip()
            if not body: continue
            for sub in re.split(r'\n(?=### )', body):
                sub = sub.strip()
                if not sub or len(sub) < 40: continue
                sub_lines = sub.split('\n')
                sub_title = sub_lines[0].lstrip('#').strip() if sub_lines[0].startswith('#') else raw_title
                sub_body  = '\n'.join(sub_lines[1:]).strip() if sub_lines[0].startswith('#') else sub
                chunks.append({
                    'module': mod, 'module_title': mod_title,
                    'source_type': 'notes', 'topic': raw_title[:120],
                    'title': sub_title[:500],
                    'content': f"Study notes — {raw_title}: {sub_title}\n\n{sub_body}",
                    'metadata': {'section': raw_title},
                })
        n = upsert_chunks(sb, chunks, oai, dry_run)
        print(f'  M{mod} notes: {n} chunks')
        total += n
    return total

# ══════════════════════════════════════════════════════════════════════
# Source 3: Flashcard Q&A
# ══════════════════════════════════════════════════════════════════════
def _extract_flashcards(html: str) -> list[tuple[str, str]]:
    """
    Robust regex extractor — finds every (question, answer) pair regardless
    of nesting depth. Strips all inner HTML tags to produce clean plain text.
    Previously the HTMLParser dropped 24 cards due to depth-tracking quirks.
    """
    def clean(s: str) -> str:
        s = re.sub(r'<[^>]+>', ' ', s)
        s = re.sub(r'&[a-z#0-9]+;', ' ', s)
        return ' '.join(s.split()).strip()

    pairs = re.findall(
        r'<div class="q">(.*?)</div>.*?<div class="a">(.*?)</div>',
        html, re.DOTALL
    )
    return [(clean(q), clean(a)) for q, a in pairs if clean(q) and clean(a)]


def ingest_flashcards(oai, sb, only_module=None, dry_run=False) -> int:
    total = 0
    all_modules = ['01','02','03','04','05','06','07','08','09']
    modules = [only_module] if only_module else all_modules
    for mod in modules:
        path = FLASHCARDS / f'module-{mod}' / 'index.html'
        if not path.exists(): continue
        cards = _extract_flashcards(path.read_text(encoding='utf-8'))
        mod_title = MODULE_TITLES.get(mod, '')
        chunks = [{
            'module': mod, 'module_title': mod_title,
            'source_type': 'flashcard', 'topic': None,
            'title': q[:500], 'content': f"Q: {q}\n\nA: {a}",
            'metadata': {},
        } for q, a in cards]
        n = upsert_chunks(sb, chunks, oai, dry_run)
        print(f'  M{mod} flashcards: {n} chunks ({len(cards)} cards)')
        total += n
    return total

# ══════════════════════════════════════════════════════════════════════
# Source 4: Reference HTML learning pages
# (learn-visually, deep-dive, cluster-explained, visual-concepts)
# ══════════════════════════════════════════════════════════════════════
class _RefPageParser(HTMLParser):
    """
    Extract section-level content from reference HTML pages.
    Handles both <section id="..."> (learn-visually, deep-dive)
    and <div class="section" id="..."> (cluster-explained).
    Skips script/style/svg/aside.
    """
    SKIP = {'script', 'style', 'svg', 'defs', 'noscript', 'aside', 'nav'}

    def __init__(self):
        super().__init__()
        self._skip_depth = 0
        self._sections: list[dict] = []
        self._cur: dict | None = None
        self._in_h2 = self._in_h3 = False
        self._div_section_depth = 0
        self._buf: list[str] = []

    def _is_section_div(self, attrs):
        cls = attrs.get('class', '')
        return 'section' in cls.split()

    def _flush_section(self):
        if self._cur:
            self._cur['content'] = re.sub(r'\s+', ' ', ' '.join(self._buf)).strip()
            self._buf = []
            if self._cur['h2'] and len(self._cur['content']) > 40:
                self._sections.append(self._cur)
        self._cur = None

    def handle_starttag(self, tag, attrs):
        if tag in self.SKIP: self._skip_depth += 1; return
        if self._skip_depth: return
        attrs = dict(attrs)
        # Both <section> and <div class="section"> open a new section block
        if tag == 'section' or (tag == 'div' and self._is_section_div(attrs)):
            if tag == 'div': self._div_section_depth += 1
            self._flush_section()
            self._cur = {'id': attrs.get('id',''), 'h2': '', 'h3s': [], 'content': ''}
            self._buf = []
        elif tag == 'h2': self._in_h2 = True
        elif tag == 'h3': self._in_h3 = True

    def handle_endtag(self, tag):
        if tag in self.SKIP and self._skip_depth:
            self._skip_depth -= 1; return
        if tag == 'section': self._flush_section()
        elif tag == 'div' and self._div_section_depth > 0:
            self._div_section_depth -= 1
            if self._div_section_depth == 0: self._flush_section()
        elif tag == 'h2': self._in_h2 = False
        elif tag == 'h3': self._in_h3 = False

    def handle_data(self, data):
        if self._skip_depth or not self._cur: return
        text = data.strip()
        if not text: return
        if self._in_h2:
            self._cur['h2'] = (self._cur['h2'] + ' ' + text).strip()
        elif self._in_h3:
            self._cur['h3s'].append(text)
            self._buf.append(text)
        else:
            self._buf.append(text)

    def handle_entityref(self, name):
        self.handle_data({'amp':'&','lt':'<','gt':'>','nbsp':' ',
                          'quot':'"','apos':"'"}.get(name,' '))
    def handle_charref(self, name):
        try:
            self.handle_data(chr(int(name[1:],16) if name.startswith('x') else int(name)))
        except Exception: pass


def _extract_visual_concepts(html: str, mod: str, page_title: str) -> list[dict]:
    """
    visual-concepts.html uses a JS TOPICS array with title/frames[].
    Extract topic → frame explain+narration as RAG chunks.
    """
    mod_title = MODULE_TITLES.get(mod, '')
    chunks = []

    # Find each top-level topic block: {title:'...' ... frames:[...]}
    # Extract topic titles
    topic_titles = re.findall(r"\{title:'([^']+)',\s*icon:", html)
    # Extract all frame explains and narrations paired with their surrounding topic
    # We walk through the file, tracking which topic we're in
    segments = re.split(r"\{title:'([^']+)',\s*icon:", html)
    # segments[0] = before first topic; then alternating title, body
    for i in range(1, len(segments), 2):
        if i + 1 >= len(segments): break
        topic_title = segments[i].strip()
        topic_body  = segments[i+1]

        # Extract all (frame_title, explain, narration) triples from this topic body
        frame_titles  = re.findall(r"title:'([^']+)'", topic_body)
        explains      = re.findall(r"explain:'([^']+)'", topic_body)
        narrations    = re.findall(r"narration:'([^']+)'", topic_body)

        for j, narr in enumerate(narrations):
            ft      = frame_titles[j] if j < len(frame_titles) else f'Frame {j+1}'
            explain = explains[j]     if j < len(explains)     else ''
            content = f"Visual Concept — {topic_title}\n\nStep: {ft}\n\n"
            if explain: content += f"Explanation: {explain}\n\n"
            content += f"Narration: {narr}"
            chunks.append({
                'module': mod, 'module_title': mod_title,
                'source_type': 'reference',
                'topic': topic_title,
                'title': ft,
                'content': content,
                'metadata': {'page': 'visual-concepts.html', 'topic': topic_title},
            })
    return chunks


def ingest_reference_pages(oai, sb, only_module=None, dry_run=False) -> int:
    total = 0
    pages = [(f, m, t) for f, m, t in REFERENCE_PAGES
             if not only_module or m == only_module]
    for fname, mod, page_title in pages:
        path = LAB_DIR / fname
        if not path.exists(): print(f'  [skip] {fname}'); continue

        html      = path.read_text(encoding='utf-8')
        mod_title = MODULE_TITLES.get(mod, '')
        chunks    = []

        # visual-concepts.html: JS-driven content → custom extractor
        if fname == 'visual-concepts.html':
            chunks = _extract_visual_concepts(html, mod, page_title)
        else:
            parser = _RefPageParser()
            parser.feed(html)
            for sec in parser._sections:
                h2      = sec.get('h2', '').strip()
                content = sec.get('content', '').strip()
                if not h2 or len(content) < 50: continue
                content = re.sub(r'  +', ' ', content)
                # Split very long sections to stay under 4000 chars
                for i, ct in enumerate([content[j:j+3500] for j in range(0,len(content),3500)]):
                    chunks.append({
                        'module': mod, 'module_title': mod_title,
                        'source_type': 'reference',
                        'topic': h2,
                        'title': f"{h2}{' (cont.)' if i else ''}",
                        'content': f"Reference — {page_title}\n\n{h2}\n\n{ct}",
                        'metadata': {
                            'page': fname, 'anchor': sec.get('id',''),
                            'h3s': sec.get('h3s',[])[:5],
                        },
                    })

        n = upsert_chunks(sb, chunks, oai, dry_run)
        print(f'  {fname}: {n} chunks')
        total += n
    return total


# ══════════════════════════════════════════════════════════════════════
# Source 5: Animated lab step narrations (engine-based labs)
# ══════════════════════════════════════════════════════════════════════
def ingest_lab_steps(oai, sb, only_module=None, dry_run=False) -> int:
    total = 0
    labs = {k:v for k,v in ENGINE_LABS.items()
            if not only_module or k == only_module}

    for mod, fname in labs.items():
        path = LAB_DIR / fname
        if not path.exists(): print(f'  [skip] {fname}'); continue

        html = path.read_text(encoding='utf-8')
        mod_title = MODULE_TITLES.get(mod, '')

        # Extract LAB_CONFIG badge/title
        config_m = re.search(r"badge:'([^']*)'.*?title:'([^']*)'", html, re.DOTALL)
        lab_title = config_m.group(2) if config_m else fname

        # Extract all step objects: {section:'...', title:'...', narration:'...'}
        # Use a broader pattern that handles multiline
        step_pattern = re.compile(
            r'\{(?:[^{}]|\{[^{}]*\})*?'
            r"section:'([^']*)'[^}]*?"
            r"title:'([^']*)'[^}]*?"
            r"narration:'([^']*)'",
            re.DOTALL
        )
        # Fallback: simpler pattern for straightforward steps
        narration_pattern = re.compile(
            r"section:'([^']+)'.*?title:'([^']+)'.*?narration:'([^']+)'",
            re.DOTALL
        )

        # Extract steps by finding each {section:... narration:...} block
        # Walk through the JS finding paired section+narration
        sections_found = re.findall(r"section:'([^']*)'", html)
        titles_found   = re.findall(r"title:'([^']*)'", html)
        narr_found     = re.findall(r"narration:'([^']*)'", html)

        chunks = []
        # Pair up sections/titles/narrations (they appear in step order)
        n_steps = len(narr_found)
        for i, narr in enumerate(narr_found):
            if not narr.strip(): continue
            section = sections_found[i] if i < len(sections_found) else ''
            title   = titles_found[i]   if i < len(titles_found)   else f'Step {i+1}'
            chunks.append({
                'module': mod, 'module_title': mod_title,
                'source_type': 'lab_step',
                'topic': section or lab_title,
                'title': title,
                'content': (
                    f"Animated Lab — {lab_title}\n"
                    f"Section: {section}\nStep: {title}\n\n{narr}"
                ),
                'metadata': {'lab': fname, 'step_index': i},
            })

        n = upsert_chunks(sb, chunks, oai, dry_run)
        print(f'  M{mod} lab steps ({fname}): {n} chunks')
        total += n
    return total

# ══════════════════════════════════════════════════════════════════════
# Source 6: Lab notes (lab-1-deploying-pods/notes/*.md)
# Maps to module 01 — K8s hands-on concepts
# ══════════════════════════════════════════════════════════════════════
def ingest_lab_notes(oai, sb, only_module=None, dry_run=False) -> int:
    if only_module and only_module != '01':
        return 0

    notes_dir = LAB_DIR / 'notes'
    if not notes_dir.exists(): return 0

    mod, mod_title = '01', MODULE_TITLES['01']
    chunks = []

    for md_path in sorted(notes_dir.glob('*.md')):
        text = md_path.read_text(encoding='utf-8')
        file_label = md_path.stem  # e.g. "01-concepts"

        # Split on ## headers
        sections = re.split(r'\n(?=## )', text)
        for sec in sections:
            sec = sec.strip()
            if not sec or len(sec) < 60: continue
            lines = sec.split('\n')
            raw_title = lines[0].lstrip('#').strip() or file_label
            body = '\n'.join(lines[1:]).strip()
            if not body or len(body) < 40: continue

            # Sub-split on ### headers for granularity
            for sub in re.split(r'\n(?=### )', body):
                sub = sub.strip()
                if not sub or len(sub) < 30: continue
                sub_lines = sub.split('\n')
                sub_title = sub_lines[0].lstrip('#').strip() if sub_lines[0].startswith('#') else raw_title
                sub_body  = '\n'.join(sub_lines[1:]).strip() if sub_lines[0].startswith('#') else sub
                if not sub_body: continue
                chunks.append({
                    'module': mod, 'module_title': mod_title,
                    'source_type': 'lab_notes',
                    'topic': raw_title[:120],
                    'title': sub_title[:500],
                    'content': f"Lab notes — {raw_title}: {sub_title}\n\n{sub_body}",
                    'metadata': {'file': md_path.name, 'section': raw_title},
                })

    n = upsert_chunks(sb, chunks, oai, dry_run)
    print(f'  M01 lab notes ({len(chunks)} sections): {n} chunks')
    return n

# ══════════════════════════════════════════════════════════════════════
# CLI
# ══════════════════════════════════════════════════════════════════════
ALL_SOURCES = ['slides', 'notes', 'flashcards', 'reference', 'lab_steps', 'lab_notes']

def main():
    ap = argparse.ArgumentParser(description='Ingest EKS course content into Supabase pgvector')
    ap.add_argument('--source',  choices=ALL_SOURCES, help='Only ingest one source type')
    ap.add_argument('--module',  help='Only ingest one module (e.g. 01)')
    ap.add_argument('--clear',   action='store_true', help='Delete all rows before ingesting')
    ap.add_argument('--dry-run', action='store_true', help='Count chunks without embedding or inserting')
    ap.add_argument('--env',     default='.env.local', help='Path to .env file')
    args = ap.parse_args()

    # Load .env
    env_path = REPO_ROOT / args.env
    if env_path.exists():
        for line in env_path.read_text().splitlines():
            line = line.strip()
            if line and not line.startswith('#') and '=' in line:
                k, _, v = line.partition('=')
                # Always override — ensures .env.local wins over stale shell env
                os.environ[k.strip()] = v.strip().strip('"').strip("'")

    oai, sb = get_clients()

    if args.clear and not args.dry_run:
        print(f'⚠ Clearing knowledge_chunks rows with corpus={CORPUS}…')
        sb.table('knowledge_chunks').delete().eq('corpus', CORPUS).execute()
        print('  cleared.')

    total = 0
    sources = [args.source] if args.source else ALL_SOURCES
    dry = args.dry_run

    for source in sources:
        print(f'\n── {source.upper()} ──')
        if source == 'slides':
            total += ingest_slides(oai, sb, args.module, dry)
        elif source == 'notes':
            total += ingest_notes(oai, sb, args.module, dry)
        elif source == 'flashcards':
            total += ingest_flashcards(oai, sb, args.module, dry)
        elif source == 'reference':
            total += ingest_reference_pages(oai, sb, args.module, dry)
        elif source == 'lab_steps':
            total += ingest_lab_steps(oai, sb, args.module, dry)
        elif source == 'lab_notes':
            total += ingest_lab_notes(oai, sb, args.module, dry)

    print(f'\n✓ {"[dry-run] " if dry else ""}Done — {total} total chunks {"counted" if dry else "upserted"}')

if __name__ == '__main__':
    main()
