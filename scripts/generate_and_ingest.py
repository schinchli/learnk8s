#!/usr/bin/env python3
#
# DEPRECATED — use scripts/ingest.py instead.
# This script writes to the legacy public.eks_knowledge table which has
# been ARCHIVED → public.eks_knowledge_archived_2026_05 in LMS migration 016.
# Running this script unmodified will error on "relation does not exist".
# To re-generate course content, use:
#   /tmp $ ~/Documents/Projects/learnk8s/.venv/bin/python \
#         ~/Documents/Projects/learnk8s/scripts/ingest.py --clear
#
"""
Universal Generate-and-Ingest Pipeline
=======================================
Drop any PDF or PPTX → generates quiz, reels, lab steps, summaries
via Claude → embeds with OpenAI → stores in Supabase pgvector.

Usage:
  python3 scripts/generate_and_ingest.py --file slides.pptx --module 01
  python3 scripts/generate_and_ingest.py --file report.pdf  --module custom --title "My Report"
  python3 scripts/generate_and_ingest.py --module 01 --types quiz,reel  # regen specific types
  python3 scripts/generate_and_ingest.py --from-corpus --module 01       # use existing corpus.json

Environment variables required:
  ANTHROPIC_API_KEY  - for content generation (Claude Haiku)
  OPENAI_API_KEY     - for embeddings (text-embedding-3-small)
  SUPABASE_URL       - Supabase project URL
  SUPABASE_SERVICE_ROLE_KEY - for writing to Supabase

Content types generated:
  quiz      - 5 MC questions per topic with correct answer + explanation
  reel      - 60-90s Instagram/TikTok narration script per topic
  lab_step  - step-by-step visual lab description per topic
  summary   - concise 3-bullet summary per topic
"""

import os, sys, json, re, time, argparse, textwrap
from pathlib import Path

REPO_ROOT   = Path(__file__).parent.parent
CORPUS_FILE = REPO_ROOT / 'scripts' / 'corpus.json'
EMBED_MODEL = 'text-embedding-3-small'
GEN_MODEL   = 'claude-haiku-4-5-20251001'

MODULE_TITLES = {
    '00': 'Course Overview',
    '01': 'Kubernetes Fundamentals',
    '02': 'Amazon EKS Fundamentals',
    '03': 'Building and Maintaining an EKS Cluster',
    '04': 'Deploying Applications to EKS',
    '05': 'Managing Applications at Scale',
    '06': 'Managing Networking in Amazon EKS',
    '07': 'Configuring Observability in Amazon EKS',
    '08': 'Managing Storage in Amazon EKS',
    '09': 'Managing Security in Amazon EKS',
    '10': 'Course Wrap-Up',
}

# ── Prompts ──────────────────────────────────────────────────────────────

QUIZ_PROMPT = """\
You are an expert AWS instructor creating exam-prep quiz questions.
Given these course slides about "{topic}", generate exactly 5 multiple-choice questions.

Slides content:
{content}

Rules:
- Questions must be answerable from the slides ONLY
- 4 answer choices (A, B, C, D) per question
- Include the correct answer letter and a 1-sentence explanation
- Vary difficulty: 2 recall, 2 application, 1 analysis
- Focus on exam-relevant facts (not logistics/slide numbers)

Output as JSON array:
[
  {{
    "q": "question text",
    "choices": {{"A":"...", "B":"...", "C":"...", "D":"..."}},
    "answer": "B",
    "explanation": "because..."
  }},
  ...
]
Output JSON only. No markdown fences."""

REEL_PROMPT = """\
You are creating a 60-second educational Instagram/TikTok reel script about "{topic}".
The audience is cloud engineers preparing for the AWS EKS course exam.

Source content:
{content}

Rules:
- Hook in first 5 words (grab attention immediately)
- 8-12 punchy sentences max
- Short sentences (max 15 words each)
- Include 1 memorable analogy
- End with a clear takeaway or call to action
- Conversational tone, not corporate
- No emojis in the script itself

Format:
HOOK: [first sentence]
SCRIPT:
[full script, one sentence per line]
VISUAL_CUE: [one-line description of what should be on screen]
HASHTAGS: [5 relevant hashtags]"""

LAB_PROMPT = """\
You are writing step-by-step visual lab instructions for "{topic}" in an animated SVG lab.
Source content:
{content}

Generate 4-6 lab steps. Each step teaches one key concept visually.

Output as JSON array:
[
  {{
    "step": 1,
    "title": "short step title",
    "narration": "1-2 sentence explanation spoken aloud as SVG animates",
    "visual": "description of what the SVG should show: boxes, arrows, labels",
    "key_fact": "the one thing the learner must remember from this step"
  }},
  ...
]
Output JSON only. No markdown fences."""

SUMMARY_PROMPT = """\
Summarise "{topic}" from this AWS EKS course content in exactly 3 bullet points.
Each bullet: bold key term + explanation in plain English (max 20 words).

Content:
{content}

Format:
• **Term**: explanation
• **Term**: explanation
• **Term**: explanation

Plus one "Exam trap" warning at the end:
⚠ Exam trap: ..."""

# ── Clients ──────────────────────────────────────────────────────────────

def get_clients():
    from anthropic import Anthropic
    from openai import OpenAI
    from supabase import create_client
    missing = [k for k in ('ANTHROPIC_API_KEY','OPENAI_API_KEY','SUPABASE_URL','SUPABASE_SERVICE_ROLE_KEY') if not os.environ.get(k)]
    if missing:
        sys.exit(f'Missing env vars: {", ".join(missing)}\nRun with keys inline:\n  ANTHROPIC_API_KEY=sk-ant-... OPENAI_API_KEY=sk-... SUPABASE_URL=... SUPABASE_SERVICE_ROLE_KEY=... python3 scripts/generate_and_ingest.py')
    return (
        Anthropic(api_key=os.environ['ANTHROPIC_API_KEY']),
        OpenAI(api_key=os.environ['OPENAI_API_KEY']),
        create_client(os.environ['SUPABASE_URL'], os.environ['SUPABASE_SERVICE_ROLE_KEY'])
    )

# ── Extraction ────────────────────────────────────────────────────────────

def extract_pptx(path: Path) -> list[dict]:
    from pptx import Presentation
    prs = Presentation(str(path))
    chunks, topic = [], 'General'
    for idx, slide in enumerate(prs.slides):
        texts = [s.text.strip().replace('\x0b','\n') for s in slide.shapes if hasattr(s,'text') and s.text.strip()]
        if not texts: continue
        if len(texts) <= 2 and all(len(t) < 90 for t in texts):
            t = texts[0].replace('\n',' ').strip()
            if t and 'running containers' not in t.lower(): topic = t
            continue
        title = texts[0].replace('\n',' ').strip()
        content = '\n'.join(texts)
        tl = title.lower()
        if 'knowledge check' in tl: topic = 'Knowledge Check'
        elif 'module summary' in tl: topic = 'Module Summary'
        chunks.append({'slide': idx+1, 'title': title[:500], 'content': content[:4000], 'topic': topic})
    return chunks

def extract_pdf(path: Path) -> list[dict]:
    try:
        import pdfplumber
    except ImportError:
        sys.exit('Install pdfplumber: pip install pdfplumber')
    chunks = []
    with pdfplumber.open(str(path)) as pdf:
        for i, page in enumerate(pdf.pages):
            text = (page.extract_text() or '').strip()
            if len(text) < 50: continue
            lines = text.split('\n')
            title = lines[0][:200] if lines else f'Page {i+1}'
            chunks.append({'slide': i+1, 'title': title, 'content': text[:4000], 'topic': f'Page {i+1}'})
    return chunks

def group_by_topic(chunks: list[dict]) -> dict[str, list[dict]]:
    groups: dict[str, list[dict]] = {}
    for c in chunks:
        t = c.get('topic') or 'General'
        if t in ('Knowledge Check', 'Module Summary', 'Lab Exercise', 'Demonstration'):
            continue  # skip non-content slides
        groups.setdefault(t, []).append(c)
    return groups

# ── Generation ────────────────────────────────────────────────────────────

def call_claude(ant_client, prompt: str, max_tokens: int = 1200) -> str:
    for attempt in range(3):
        try:
            msg = ant_client.messages.create(
                model=GEN_MODEL, max_tokens=max_tokens,
                messages=[{'role':'user','content':prompt}]
            )
            return msg.content[0].text
        except Exception as e:
            if attempt == 2: raise
            print(f'    retrying Claude ({e})...')
            time.sleep(2**attempt)
    return ''

def gen_quiz(ant, topic: str, content: str) -> list[dict]:
    prompt = QUIZ_PROMPT.format(topic=topic, content=content[:3000])
    raw = call_claude(ant, prompt, 1400)
    try:
        return json.loads(raw)
    except json.JSONDecodeError:
        m = re.search(r'\[.*\]', raw, re.DOTALL)
        if m:
            try: return json.loads(m.group())
            except: pass
    return []

def gen_reel(ant, topic: str, content: str) -> str:
    return call_claude(ant, REEL_PROMPT.format(topic=topic, content=content[:3000]), 600)

def gen_lab_steps(ant, topic: str, content: str) -> list[dict]:
    prompt = LAB_PROMPT.format(topic=topic, content=content[:3000])
    raw = call_claude(ant, prompt, 1000)
    try:
        return json.loads(raw)
    except json.JSONDecodeError:
        m = re.search(r'\[.*\]', raw, re.DOTALL)
        if m:
            try: return json.loads(m.group())
            except: pass
    return []

def gen_summary(ant, topic: str, content: str) -> str:
    return call_claude(ant, SUMMARY_PROMPT.format(topic=topic, content=content[:3000]), 400)

# ── Embedding + storage ───────────────────────────────────────────────────

def embed_batch(oai, texts: list[str]) -> list[list[float]]:
    results = []
    for i in range(0, len(texts), 50):
        batch = [t[:8000] for t in texts[i:i+50]]
        for attempt in range(3):
            try:
                resp = oai.embeddings.create(model=EMBED_MODEL, input=batch)
                results.extend([d.embedding for d in resp.data])
                time.sleep(0.15)
                break
            except Exception as e:
                if attempt == 2: raise
                time.sleep(2**attempt)
    return results

def upsert_rows(sb, rows: list[dict]):
    for i in range(0, len(rows), 50):
        sb.table('eks_knowledge').insert(rows[i:i+50]).execute()

# ── Content-type builders ─────────────────────────────────────────────────

def build_quiz_chunks(mod, mod_title, topic, qlist):
    chunks = []
    for q in qlist:
        choices = q.get('choices', {})
        choices_str = '  '.join(f"{k}. {v}" for k,v in choices.items())
        content = (
            f"Quiz — {topic}\n\n"
            f"Q: {q.get('q','')}\n\n"
            f"Choices: {choices_str}\n\n"
            f"Answer: {q.get('answer','')} — {q.get('explanation','')}"
        )
        chunks.append({'module':mod,'module_title':mod_title,'source_type':'quiz',
                       'topic':topic,'title':q.get('q','')[:500],
                       'content':content[:4000],'metadata':{'answer':q.get('answer','')}})
    return chunks

def build_reel_chunk(mod, mod_title, topic, script: str):
    return {'module':mod,'module_title':mod_title,'source_type':'reel',
            'topic':topic,'title':f'Reel: {topic}'[:500],
            'content':f'Reel script — {topic}\n\n{script}'[:4000],'metadata':{}}

def build_lab_chunks(mod, mod_title, topic, steps):
    chunks = []
    for st in steps:
        content = (
            f"Lab step {st.get('step','')} — {topic}: {st.get('title','')}\n\n"
            f"Narration: {st.get('narration','')}\n\n"
            f"Visual: {st.get('visual','')}\n\n"
            f"Key fact: {st.get('key_fact','')}"
        )
        chunks.append({'module':mod,'module_title':mod_title,'source_type':'lab_step',
                       'topic':topic,'title':f"Step {st.get('step','')} — {st.get('title',''[:200])}",
                       'content':content[:4000],'metadata':{'step':st.get('step',0)}})
    return chunks

def build_summary_chunk(mod, mod_title, topic, text: str):
    return {'module':mod,'module_title':mod_title,'source_type':'summary',
            'topic':topic,'title':f'Summary: {topic}'[:500],
            'content':f'Summary — {topic}\n\n{text}'[:4000],'metadata':{}}

# ── Main ──────────────────────────────────────────────────────────────────

def main():
    ap = argparse.ArgumentParser(description='Generate quiz/reels/lab from PPTX or PDF + store in Supabase')
    ap.add_argument('--file',         help='Path to PPTX or PDF file')
    ap.add_argument('--module',       default='custom', help='Module number, e.g. 01')
    ap.add_argument('--title',        help='Module title (auto-detected if known module)')
    ap.add_argument('--from-corpus',  action='store_true', help='Read from scripts/corpus.json instead of file')
    ap.add_argument('--types',        default='quiz,reel,lab_step,summary',
                                      help='Comma-separated content types to generate')
    ap.add_argument('--topics',       help='Only generate for specific topics (comma-separated)')
    ap.add_argument('--dry-run',      action='store_true', help='Show what would be generated, no API calls')
    ap.add_argument('--env',          default='.env.local', help='Path to .env file')
    args = ap.parse_args()

    # Load .env
    env_path = REPO_ROOT / args.env
    if env_path.exists():
        for line in env_path.read_text().splitlines():
            line = line.strip()
            if line and not line.startswith('#') and '=' in line:
                k, _, v = line.partition('=')
                # Always override — ensures .env.local wins over stale shell env
                os.environ[k.strip()] = v.strip().strip('"').strip("'")

    mod        = args.module
    mod_title  = args.title or MODULE_TITLES.get(mod, f'Module {mod}')
    types      = [t.strip() for t in args.types.split(',')]
    only_topics = [t.strip() for t in args.topics.split(',')] if args.topics else None

    # ── Extract source chunks ──────────────────────────────────────
    if args.from_corpus:
        corpus = json.loads(CORPUS_FILE.read_text())
        raw_chunks = [c for c in corpus if c['module'] == mod and c['source_type'] == 'slide']
    elif args.file:
        p = Path(args.file)
        if not p.exists(): sys.exit(f'File not found: {p}')
        raw_chunks = extract_pptx(p) if p.suffix.lower() in ('.pptx','.ppt') else extract_pdf(p)
    else:
        sys.exit('Provide --file or --from-corpus')

    topic_groups = group_by_topic(raw_chunks)
    if only_topics:
        topic_groups = {k:v for k,v in topic_groups.items() if k in only_topics}

    print(f'\n{"─"*60}')
    print(f'  Module {mod}: {mod_title}')
    print(f'  Source: {len(raw_chunks)} slides  →  {len(topic_groups)} topics')
    print(f'  Types:  {", ".join(types)}')
    print(f'  Topics: {list(topic_groups.keys())[:6]}...' if len(topic_groups)>6 else f'  Topics: {list(topic_groups.keys())}')
    print(f'{"─"*60}\n')

    if args.dry_run:
        total = sum(1 for _ in topic_groups for _ in types)
        print(f'  --dry-run: would generate ~{total} content items across {len(topic_groups)} topics')
        return

    ant, oai, sb = get_clients()

    all_chunks = []
    total_topics = len(topic_groups)

    for ti, (topic, slides) in enumerate(topic_groups.items(), 1):
        print(f'  [{ti}/{total_topics}] {topic[:50]}')
        content = '\n\n'.join(f"[Slide {s['slide']}] {s['title']}\n{s['content']}" for s in slides)

        if 'quiz' in types:
            print(f'    → quiz...', end='', flush=True)
            qlist = gen_quiz(ant, topic, content)
            chunks = build_quiz_chunks(mod, mod_title, topic, qlist)
            all_chunks.extend(chunks)
            print(f' {len(chunks)} questions')

        if 'reel' in types:
            print(f'    → reel...', end='', flush=True)
            script = gen_reel(ant, topic, content)
            all_chunks.append(build_reel_chunk(mod, mod_title, topic, script))
            print(' done')

        if 'lab_step' in types:
            print(f'    → lab steps...', end='', flush=True)
            steps = gen_lab_steps(ant, topic, content)
            chunks = build_lab_chunks(mod, mod_title, topic, steps)
            all_chunks.extend(chunks)
            print(f' {len(chunks)} steps')

        if 'summary' in types:
            print(f'    → summary...', end='', flush=True)
            text = gen_summary(ant, topic, content)
            all_chunks.append(build_summary_chunk(mod, mod_title, topic, text))
            print(' done')

        time.sleep(0.3)  # rate limit courtesy

    print(f'\n  Generated {len(all_chunks)} total chunks. Embedding...')
    texts = [c['content'] for c in all_chunks]
    embeddings = embed_batch(oai, texts)

    rows = [
        {**c, 'embedding': emb, 'metadata': c.get('metadata',{})}
        for c, emb in zip(all_chunks, embeddings)
    ]

    print(f'  Storing {len(rows)} rows in Supabase...')
    upsert_rows(sb, rows)

    # ── Save generated content as JSON for reference ───────────────
    out = REPO_ROOT / 'scripts' / f'generated_module_{mod}.json'
    out.write_text(json.dumps(all_chunks, indent=2, ensure_ascii=False))

    print(f'\n✓  Done — {len(rows)} chunks stored')
    print(f'   Saved: {out}')
    print(f'\n   Content types added to eks_knowledge:')
    from collections import Counter
    for typ, cnt in Counter(c["source_type"] for c in all_chunks).items():
        print(f'   {typ:<12} {cnt:>4}')

if __name__ == '__main__':
    main()
