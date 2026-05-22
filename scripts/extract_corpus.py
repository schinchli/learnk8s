#!/usr/bin/env python3
"""
Stage 1 — pure extraction, zero API calls, zero credentials.

Reads every PPTX + notes.md + flashcard HTML and writes
scripts/corpus.json ready for the embed stage.

Usage:
  python3 scripts/extract_corpus.py
"""

import json, re, sys
from pathlib import Path
from html.parser import HTMLParser
from pptx import Presentation

REPO_ROOT  = Path(__file__).parent.parent
PPTX_BASE  = Path('/Users/schinchli/Downloads/200-COREKS-22-EN-PPTX.2.2.3-20260417145107 2/')
FLASHCARDS = REPO_ROOT / 'course-flashcards'
OUT_FILE   = REPO_ROOT / 'scripts' / 'corpus.json'

MODULE_META = {
    '00': ('Course Overview',                              '00_CourseOverview_InstructorDeck.pptx'),
    '01': ('Kubernetes Fundamentals',                      '01_KubernetesFundamentals_InstructorDeck.pptx'),
    '02': ('Amazon EKS Fundamentals',                      '02_AmazonEKSFundamentals_InstructorDeck.pptx'),
    '03': ('Building and Maintaining an EKS Cluster',      '03_BuildingandMaintaininganAmazonEKSCluster_InstructorDeck.pptx'),
    '04': ('Deploying Applications to EKS',                '04_DeployingApplicationstoyourEKSCluster_InstructorDeck.pptx'),
    '05': ('Managing Applications at Scale',               '05_ManagingApplicationsatScaleinAmazonEKS_InstructorDeck.pptx'),
    '06': ('Managing Networking in Amazon EKS',            '06_ManagingNetworkinginAmazonEKS_InstructorDeck.pptx'),
    '07': ('Configuring Observability in Amazon EKS',      '07_ConfiguringObservabilityinanAmazonEKSCluster_InstructorDeck.pptx'),
    '08': ('Managing Storage in Amazon EKS',               '08_ManagingStorageinAmazonEKS_InstructorDeck.pptx'),
    '09': ('Managing Security in Amazon EKS',              '09_ManagingSecurityinAmazonEKS_InstructorDeck.pptx'),
    '10': ('Course Wrap-Up',                               '10_CourseWrapUp_InstructorDeck.pptx'),
}

# ── helpers ──────────────────────────────────────────────────────────

def clean(text: str) -> str:
    """Normalise whitespace and remove vertical-tab chars used in PPTX."""
    return re.sub(r'\s+', ' ', text.replace('\x0b', '\n').replace('\r', '')).strip()


def slide_chunks(mod: str, mod_title: str, pptx_path: Path) -> list[dict]:
    prs     = Presentation(str(pptx_path))
    chunks  = []
    topic   = mod_title

    for idx, slide in enumerate(prs.slides):
        texts = [
            clean(s.text)
            for s in slide.shapes
            if hasattr(s, 'text') and clean(s.text)
        ]
        if not texts:
            continue

        # Section-divider slides: short, update running topic, skip as chunk
        if len(texts) <= 2 and all(len(t) < 90 for t in texts):
            candidate = texts[0].replace('\n', ' ').strip()
            if candidate and 'running containers' not in candidate.lower():
                topic = candidate
            continue

        slide_num = idx + 1
        title     = texts[0].replace('\n', ' ').strip()
        body      = '\n'.join(texts[1:])

        # Update topic from content hints
        tl = title.lower()
        if 'knowledge check' in tl:
            topic = 'Knowledge Check'
        elif 'module summary' in tl:
            topic = 'Module Summary'
        elif tl.startswith('lab '):
            topic = 'Lab Exercise'
        elif 'demonstration' in tl:
            topic = 'Demonstration'

        content = f"{title}\n\n{body}".strip()
        if len(content) < 30:
            continue

        chunks.append({
            'module':       mod,
            'module_title': mod_title,
            'source_type':  'slide',
            'topic':        topic,
            'title':        title[:500],
            'content':      content[:4000],
            'metadata':     {'slide': slide_num, 'pptx': pptx_path.name},
        })

    return chunks


def notes_chunks(mod: str, mod_title: str) -> list[dict]:
    path = FLASHCARDS / f'module-{mod}' / 'notes.md'
    if not path.exists():
        return []

    text   = path.read_text(encoding='utf-8')
    chunks = []

    for section in re.split(r'\n(?=## )', text):
        section = section.strip()
        if not section or len(section) < 60:
            continue
        lines     = section.split('\n')
        h2_title  = lines[0].lstrip('#').strip()
        body      = '\n'.join(lines[1:]).strip()
        if not body:
            continue

        # Split long sections at ### boundaries
        for sub in re.split(r'\n(?=### )', body):
            sub = sub.strip()
            if len(sub) < 40:
                continue
            sub_lines  = sub.split('\n')
            sub_title  = sub_lines[0].lstrip('#').strip() if sub_lines[0].startswith('#') else h2_title
            sub_body   = '\n'.join(sub_lines[1:]).strip() if sub_lines[0].startswith('#') else sub

            content = f"Study notes — {h2_title}: {sub_title}\n\n{sub_body}"
            chunks.append({
                'module':       mod,
                'module_title': mod_title,
                'source_type':  'notes',
                'topic':        h2_title[:120],
                'title':        sub_title[:500],
                'content':      content[:4000],
                'metadata':     {'section': h2_title},
            })

    return chunks


class _CardParser(HTMLParser):
    def __init__(self):
        super().__init__()
        self._in_q = self._in_a = False
        self._dq = self._da = 0
        self._q = self._a = ''
        self._stack: list[tuple[str, str]] = []
        self.cards: list[tuple[str, str]] = []

    def handle_starttag(self, tag, attrs):
        cls = dict(attrs).get('class', '')
        self._stack.append((tag, cls))
        if tag == 'div' and cls == 'q':
            self._in_q = True;  self._dq = len(self._stack)
        elif tag == 'div' and cls == 'a':
            self._in_a = True;  self._da = len(self._stack)

    def handle_endtag(self, tag):
        if self._stack: self._stack.pop()
        if self._in_q and len(self._stack) < self._dq:
            self._in_q = False
        if self._in_a and len(self._stack) < self._da:
            self._in_a = False
            q = re.sub(r'\s+', ' ', self._q).strip()
            a = re.sub(r'&[a-z]+;|\s+', lambda m: ' ' if m.group().startswith('&') or m.group() in (' ', '\t') else '\n', self._a).strip()
            if q and a:
                self.cards.append((q, a))
            self._q = self._a = ''

    def handle_data(self, data):
        if self._in_q: self._q += data
        elif self._in_a: self._a += data


def flashcard_chunks(mod: str, mod_title: str) -> list[dict]:
    path = FLASHCARDS / f'module-{mod}' / 'index.html'
    if not path.exists():
        return []

    parser = _CardParser()
    parser.feed(path.read_text(encoding='utf-8'))

    chunks = []
    for q, a in parser.cards:
        q = re.sub(r'\s+', ' ', q).strip()
        a = re.sub(r'\s+', ' ', a).strip()
        if not q or not a or len(q) < 10:
            continue
        chunks.append({
            'module':       mod,
            'module_title': mod_title,
            'source_type':  'flashcard',
            'topic':        None,
            'title':        q[:500],
            'content':      f"Q: {q}\n\nA: {a}",
            'metadata':     {},
        })

    return chunks


# ── main ─────────────────────────────────────────────────────────────

def main():
    all_chunks: list[dict] = []
    stats: dict[str, dict] = {}

    for mod, (mod_title, fname) in MODULE_META.items():
        pptx_path = PPTX_BASE / fname
        s = stats[mod] = {'title': mod_title, 'slides': 0, 'notes': 0, 'flashcards': 0}

        # Slides
        if pptx_path.exists():
            sc = slide_chunks(mod, mod_title, pptx_path)
            all_chunks.extend(sc)
            s['slides'] = len(sc)
        else:
            print(f'  [warn] module {mod}: PPTX not found at {pptx_path}', file=sys.stderr)

        # Notes
        nc = notes_chunks(mod, mod_title)
        all_chunks.extend(nc)
        s['notes'] = len(nc)

        # Flashcards
        fc = flashcard_chunks(mod, mod_title)
        all_chunks.extend(fc)
        s['flashcards'] = len(fc)

    # Write corpus
    OUT_FILE.parent.mkdir(parents=True, exist_ok=True)
    OUT_FILE.write_text(json.dumps(all_chunks, indent=2, ensure_ascii=False), encoding='utf-8')

    # Print summary table
    total = len(all_chunks)
    print(f'\n{"Module":<6} {"Title":<48} {"Slides":>6} {"Notes":>5} {"Cards":>5} {"Total":>6}')
    print('─' * 80)
    for mod, s in stats.items():
        row_total = s['slides'] + s['notes'] + s['flashcards']
        print(f'  {mod}   {s["title"]:<48} {s["slides"]:>6} {s["notes"]:>5} {s["flashcards"]:>5} {row_total:>6}')
    print('─' * 80)
    print(f'{"TOTAL":<54} {sum(s["slides"] for s in stats.values()):>6} '
          f'{sum(s["notes"] for s in stats.values()):>5} '
          f'{sum(s["flashcards"] for s in stats.values()):>5} '
          f'{total:>6}')
    print(f'\n✓ Corpus written → {OUT_FILE}  ({total} chunks, {OUT_FILE.stat().st_size // 1024} KB)\n')


if __name__ == '__main__':
    main()
